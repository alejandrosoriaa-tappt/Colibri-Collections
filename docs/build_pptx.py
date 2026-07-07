from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta Kollybry ──────────────────────────────────────────────
GREEN_DARK   = RGBColor(0x00, 0x6C, 0x4A)   # primary #006C4A
GREEN_MED    = RGBColor(0x1D, 0xA3, 0x74)   # 500 #1da374
GREEN_LIGHT  = RGBColor(0xD6, 0xF5, 0xE8)   # 100 #d6f5e8
GREEN_XLIGHT = RGBColor(0xF0, 0xFA, 0xF5)   # 50  #f0faf5
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
ALMOST_WHITE = RGBColor(0xF8, 0xFB, 0xF9)
DARK_TEXT    = RGBColor(0x00, 0x21, 0x15)   # on-primary-container
GRAY_TEXT    = RGBColor(0x4D, 0x63, 0x59)   # secondary
ACCENT_TEAL  = RGBColor(0x3B, 0x64, 0x71)   # tertiary

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

# ── helpers ─────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=14, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None, tag=None):
    """Green top bar with title."""
    add_rect(slide, 0, 0, W, Inches(1.25), GREEN_DARK)
    add_text(slide, title,
             Inches(0.45), Inches(0.17), Inches(12), Inches(0.72),
             size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.45), Inches(0.82), Inches(10), Inches(0.38),
                 size=14, color=GREEN_LIGHT)
    if tag:
        add_rect(slide, Inches(11.5), Inches(0.35), Inches(1.6), Inches(0.45), GREEN_MED)
        add_text(slide, tag,
                 Inches(11.52), Inches(0.35), Inches(1.58), Inches(0.45),
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def bg(slide):
    add_rect(slide, 0, 0, W, H, ALMOST_WHITE)

def green_chip(slide, text, x, y, w=Inches(1.8), h=Inches(0.38)):
    add_rect(slide, x, y, w, h, GREEN_LIGHT)
    add_text(slide, text, x, y+Pt(2), w, h,
             size=11, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

def card(slide, x, y, w, h, title, bullets, title_color=GREEN_DARK):
    add_rect(slide, x, y, w, h, WHITE)
    # left accent bar
    add_rect(slide, x, y, Inches(0.06), h, GREEN_MED)
    add_text(slide, title, x+Inches(0.12), y+Inches(0.12), w-Inches(0.2), Inches(0.4),
             size=13, bold=True, color=title_color)
    txb = slide.shapes.add_textbox(x+Inches(0.12), y+Inches(0.52),
                                    w-Inches(0.2), h-Inches(0.62))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = b
        run.font.size  = Pt(11)
        run.font.color.rgb = GRAY_TEXT

def stat_box(slide, x, y, number, label, sublabel=""):
    add_rect(slide, x, y, Inches(2.6), Inches(1.6), GREEN_DARK)
    add_text(slide, number, x, y+Inches(0.15), Inches(2.6), Inches(0.75),
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y+Inches(0.85), Inches(2.6), Inches(0.4),
             size=12, bold=True, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)
    if sublabel:
        add_text(slide, sublabel, x, y+Inches(1.2), Inches(2.6), Inches(0.35),
                 size=10, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Portada
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
add_rect(s, 0, 0, W, Inches(7.5), GREEN_DARK)
# diagonal accent
add_rect(s, Inches(7.5), 0, Inches(6), H, GREEN_MED)

add_text(s, "KOLLYBRY", Inches(0.7), Inches(1.6), Inches(8), Inches(1.2),
         size=60, bold=True, color=WHITE)
add_text(s, "Comunicación masiva por WhatsApp",
         Inches(0.7), Inches(2.75), Inches(8), Inches(0.6),
         size=26, color=GREEN_LIGHT)
add_text(s, "para colegios, condominios y gimnasios",
         Inches(0.7), Inches(3.3), Inches(8), Inches(0.5),
         size=22, color=GREEN_LIGHT, italic=True)

add_text(s, "El mensaje llega donde la gente ya está,\nen menos de 3 minutos, sin instalar nada.",
         Inches(0.7), Inches(4.4), Inches(8), Inches(0.9),
         size=16, color=RGBColor(0xAE, 0xEB, 0xD2))

add_text(s, "nkuvo.com / kollybry", Inches(0.7), Inches(6.6), Inches(4), Inches(0.4),
         size=12, color=RGBColor(0x77, 0xD9, 0xB5))

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Por qué WhatsApp (benchmarks)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "¿Por qué WhatsApp?", "El canal con mayor tasa de apertura del mundo")

# stat boxes
stat_box(s, Inches(0.45), Inches(1.6),  "95–98%", "Tasa de apertura", "vs 20–30% del correo")
stat_box(s, Inches(3.3),  Inches(1.6),  "< 3 min", "Tiempo al abrirlo", "vs 6–12 h del correo")
stat_box(s, Inches(6.15), Inches(1.6),  "100M+", "Usuarios en México", "WhatsApp activos")
stat_box(s, Inches(9.0),  Inches(1.6),  "Cero", "Instalación requerida", "canal que ya usan")

add_rect(s, Inches(0.45), Inches(3.5), Inches(12.43), Inches(0.04), GREEN_LIGHT)

# comparison table header
headers = ["Canal", "Apertura", "Tiempo", "Instalar app", "Ya lo usan"]
vals    = [
    ["WhatsApp ✅",       "95–98%",   "< 3 min",  "No",     "Sí, diario"],
    ["Correo electrónico","20–30%",   "6–12 h",   "No",     "A veces"],
    ["App nativa",        "20–40%",   "variable", "Sí ❌",  "Casi nunca"],
    ["Grupos WhatsApp",   "Alta",     "inmediato","No",      "Sí, pero caótico ❌"],
]
col_x = [Inches(0.45), Inches(2.7), Inches(5.0), Inches(7.0), Inches(9.3), Inches(11.2)]
col_w = [Inches(2.2),  Inches(2.2), Inches(1.9), Inches(2.2), Inches(1.8)]
row_y = Inches(3.65)

for ci, h in enumerate(headers):
    add_rect(s, col_x[ci], row_y, col_w[ci], Inches(0.38), GREEN_DARK)
    add_text(s, h, col_x[ci]+Inches(0.05), row_y, col_w[ci], Inches(0.38),
             size=11, bold=True, color=WHITE)

for ri, row in enumerate(vals):
    y = row_y + Inches(0.38) + ri * Inches(0.52)
    fill = GREEN_XLIGHT if ri % 2 == 0 else WHITE
    for ci, cell in enumerate(row):
        add_rect(s, col_x[ci], y, col_w[ci], Inches(0.5), fill)
        col = GREEN_DARK if (ri == 0 and ci == 0) else GRAY_TEXT
        add_text(s, cell, col_x[ci]+Inches(0.05), y+Pt(2), col_w[ci]-Inches(0.1), Inches(0.46),
                 size=10, color=col, bold=(ri==0))

add_text(s, "Fuentes: Meta Business, Gartner 2024, Klaviyo Email Benchmarks 2024",
         Inches(0.45), Inches(6.9), Inches(10), Inches(0.3),
         size=9, italic=True, color=GRAY_TEXT)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Los 5 pilares de Kollybry
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Los 5 pilares de Kollybry")

pilares = [
    ("1. Canal correcto",       "WhatsApp — tasa de apertura 95–98%, sin fricción de adopción para el receptor."),
    ("2. Personalización real",  "Cada mensaje lleva nombre, monto y concepto específico de esa familia/unidad/socio. No es un blast genérico."),
    ("3. Cero fricción operativa","El admin sube un Excel, la IA interpreta columnas y salones. No requiere conocimientos técnicos."),
    ("4. Templates Meta",        "Los mensajes llegan aunque el receptor nunca haya escrito. Kollybry es outbound — diseñado exactamente para eso."),
    ("5. Memoria institucional", "Una vez configurados los mensajes de cobranza, Kollybry los recuerda para el siguiente mes. El proceso es 1 clic."),
]

for i, (title, body) in enumerate(pilares):
    row = i // 3
    col = i % 3
    if i == 3:
        x = Inches(0.45) + col * Inches(4.4)
    elif i == 4:
        x = Inches(0.45) + 1 * Inches(4.4)
    else:
        x = Inches(0.45) + col * Inches(4.4)
    y = Inches(1.55) + row * Inches(2.6)
    w = Inches(4.1)
    h = Inches(2.3)
    if i >= 3:
        y = Inches(4.1)

    add_rect(s, x, y, w, h, WHITE)
    add_rect(s, x, y, w, Inches(0.08), GREEN_MED)
    # number accent
    add_rect(s, x, y+Inches(0.08), Inches(0.5), h-Inches(0.08), GREEN_XLIGHT)
    add_text(s, str(i+1), x, y+Inches(0.35), Inches(0.5), Inches(0.5),
             size=22, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_text(s, title, x+Inches(0.55), y+Inches(0.15), w-Inches(0.65), Inches(0.4),
             size=13, bold=True, color=GREEN_DARK)
    add_text(s, body,  x+Inches(0.55), y+Inches(0.55), w-Inches(0.65), h-Inches(0.65),
             size=11, color=GRAY_TEXT)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Colegios
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Kollybry para Colegios", "Directorio, cobranza y comunicados — sin grupos de WhatsApp caóticos", tag="COLEGIOS")

# left col — problema
add_rect(s, Inches(0.45), Inches(1.45), Inches(3.9), Inches(5.7), WHITE)
add_rect(s, Inches(0.45), Inches(1.45), Inches(0.08), Inches(5.7), RGBColor(0xBA,0x1A,0x1A))
add_text(s, "❌  Sin Kollybry",
         Inches(0.6), Inches(1.55), Inches(3.7), Inches(0.4),
         size=13, bold=True, color=RGBColor(0xBA,0x1A,0x1A))
problemas = [
    "Grupos de WhatsApp por salón → caos",
    "Maestros exponen su celular personal",
    "Correos con <30% de apertura",
    "App propia que nadie instala",
    "Circulares en papel perdidas",
    "Llamadas manuales para cobrar",
    "Deuda acumulada por falta de recordatorios",
]
txb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(3.65), Inches(5.0))
tf = txb.text_frame; tf.word_wrap = True
first = True
for p_text in problemas:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(4)
    run = p.add_run(); run.text = "• " + p_text
    run.font.size = Pt(11); run.font.color.rgb = GRAY_TEXT

# right col — solución
add_rect(s, Inches(4.6), Inches(1.45), Inches(8.28), Inches(5.7), WHITE)
add_rect(s, Inches(4.6), Inches(1.45), Inches(0.08), Inches(5.7), GREEN_MED)
add_text(s, "✅  Con Kollybry",
         Inches(4.75), Inches(1.55), Inches(8.0), Inches(0.4),
         size=13, bold=True, color=GREEN_DARK)

casos = [
    ("🟢  Cobranza con semáforo de colores",
     "Sube el Excel, la IA detecta quién debe. 3 campañas listas (verde/naranja/rojo) con monto personalizado y link de pago. Enviadas en 2 minutos."),
    ("📢  Comunicados generales",
     "Suspensión de clases, eventos, emergencias. Un mensaje → todos los padres en < 3 min. Segmentado por salón o grado."),
    ("🤝  AI Import de directorio",
     "Sube tu Excel nativo (multi-hoja). La IA detecta columnas, salones y secciones. Directorio completo importado con jerarquía familiar."),
    ("✅  Confirmación de pago",
     "Al registrar un pago, mensaje automático al papá. Reduce llamadas de '¿ya les llegó?'"),
]
y = Inches(2.0)
for icon_title, body in casos:
    add_rect(s, Inches(4.68), y, Inches(8.1), Inches(1.2), GREEN_XLIGHT)
    add_text(s, icon_title, Inches(4.8), y+Inches(0.08), Inches(8.0), Inches(0.38),
             size=12, bold=True, color=GREEN_DARK)
    add_text(s, body, Inches(4.8), y+Inches(0.42), Inches(8.0), Inches(0.7),
             size=10, color=GRAY_TEXT)
    y += Inches(1.32)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Condominios
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Kollybry para Condominios", "Cuotas, avisos urgentes y asambleas — sin grupos caóticos de residentes", tag="CONDOMINIOS")

add_rect(s, Inches(0.45), Inches(1.45), Inches(3.9), Inches(5.7), WHITE)
add_rect(s, Inches(0.45), Inches(1.45), Inches(0.08), Inches(5.7), RGBColor(0xBA,0x1A,0x1A))
add_text(s, "❌  Sin Kollybry",
         Inches(0.6), Inches(1.55), Inches(3.7), Inches(0.4),
         size=13, bold=True, color=RGBColor(0xBA,0x1A,0x1A))
txb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(3.65), Inches(5.0))
tf = txb.text_frame; tf.word_wrap = True
first = True
for p_text in [
    "Grupo de WhatsApp: peleas, quejas, ruido",
    "Admin expone su número personal",
    "Correos que nadie lee",
    "Volantes físicos ignorados o perdidos",
    "Sin registro de quién recibió el aviso",
    "Morosos difíciles de contactar",
    "Asambleas con baja asistencia por falta de aviso",
]:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(4)
    run = p.add_run(); run.text = "• " + p_text
    run.font.size = Pt(11); run.font.color.rgb = GRAY_TEXT

add_rect(s, Inches(4.6), Inches(1.45), Inches(8.28), Inches(5.7), WHITE)
add_rect(s, Inches(4.6), Inches(1.45), Inches(0.08), Inches(5.7), GREEN_MED)
add_text(s, "✅  Con Kollybry",
         Inches(4.75), Inches(1.55), Inches(8.0), Inches(0.4),
         size=13, bold=True, color=GREEN_DARK)

casos_c = [
    ("🏦  Cobro de cuota de mantenimiento",
     "Reporte Excel → IA detecta quién debe → mensaje personalizado por unidad con monto y link de pago SPEI/tarjeta."),
    ("🚨  Avisos urgentes",
     "Corte de agua, falla de elevador, visita de fumigación. Mensaje masivo → todos los residentes enterados antes de llegar a casa."),
    ("📅  Convocatoria a asamblea",
     "Invitación + recordatorio automático 24h antes. Mayor asistencia sin esfuerzo del administrador."),
    ("🔐  Comunicados de seguridad",
     "Alerta de incidente, instrucciones de emergencia, acceso de visitantes. Llegada garantizada en < 3 min."),
]
y = Inches(2.0)
for icon_title, body in casos_c:
    add_rect(s, Inches(4.68), y, Inches(8.1), Inches(1.2), GREEN_XLIGHT)
    add_text(s, icon_title, Inches(4.8), y+Inches(0.08), Inches(8.0), Inches(0.38),
             size=12, bold=True, color=GREEN_DARK)
    add_text(s, body, Inches(4.8), y+Inches(0.42), Inches(8.0), Inches(0.7),
             size=10, color=GRAY_TEXT)
    y += Inches(1.32)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Gimnasios
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Kollybry para Gimnasios", "Renovaciones, campañas de reactivación y comunicados operativos", tag="GIMNASIOS")

add_rect(s, Inches(0.45), Inches(1.45), Inches(3.9), Inches(5.7), WHITE)
add_rect(s, Inches(0.45), Inches(1.45), Inches(0.08), Inches(5.7), RGBColor(0xBA,0x1A,0x1A))
add_text(s, "❌  Sin Kollybry",
         Inches(0.6), Inches(1.55), Inches(3.7), Inches(0.4),
         size=13, bold=True, color=RGBColor(0xBA,0x1A,0x1A))
txb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(3.65), Inches(5.0))
tf = txb.text_frame; tf.word_wrap = True
first = True
for p_text in [
    "WhatsApp desde el celular del dueño — no escala",
    "Correos con <30% de apertura",
    "Membresías vencidas no cobradas",
    "Sin canal para promociones masivas",
    "Socios inactivos que nunca regresan",
    "Cambios de horario que nadie ve",
]:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(4)
    run = p.add_run(); run.text = "• " + p_text
    run.font.size = Pt(11); run.font.color.rgb = GRAY_TEXT

add_rect(s, Inches(4.6), Inches(1.45), Inches(8.28), Inches(5.7), WHITE)
add_rect(s, Inches(4.6), Inches(1.45), Inches(0.08), Inches(5.7), GREEN_MED)
add_text(s, "✅  Con Kollybry",
         Inches(4.75), Inches(1.55), Inches(8.0), Inches(0.4),
         size=13, bold=True, color=GREEN_DARK)

casos_g = [
    ("💳  Renovación de membresía",
     "7 días antes del vencimiento: recordatorio con link de pago. Día del vencimiento: segundo aviso. 3 días después: oferta de recuperación."),
    ("🔄  Reactivación de socios inactivos",
     '"Te extrañamos" — segmentado por socios sin asistencia en 30+ días. Tasa de retorno hasta 3x vs correo.'),
    ("📣  Comunicados operativos",
     "Cambio de horarios, nuevo instructor, clase especial, cierre por mantenimiento. Todos los socios enterados en minutos."),
    ("🎯  Lanzamiento de servicios",
     "Nuevo servicio, reto de 21 días, evento especial. Campaña masiva con imagen adjunta para mayor impacto visual."),
]
y = Inches(2.0)
for icon_title, body in casos_g:
    add_rect(s, Inches(4.68), y, Inches(8.1), Inches(1.2), GREEN_XLIGHT)
    add_text(s, icon_title, Inches(4.8), y+Inches(0.08), Inches(8.0), Inches(0.38),
             size=12, bold=True, color=GREEN_DARK)
    add_text(s, body, Inches(4.8), y+Inches(0.42), Inches(8.0), Inches(0.7),
             size=10, color=GRAY_TEXT)
    y += Inches(1.32)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Flujo de cobranza con IA
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Flujo de cobranza con IA", "Del Excel al WhatsApp personalizado en 3 clics")

steps = [
    ("📂", "Sube\nel Excel",        "Reporte .xlsx con semáforo de colores o columna de estatus"),
    ("🤖", "IA analiza\ncolumnas",   "Claude detecta alumno, monto, estatus. Clasifica colores: verde/naranja/rojo"),
    ("👨‍👩‍👧", "Resuelve\nfamilias",   "Normaliza nombres → busca en directorio → trae papás con teléfono"),
    ("⚙️", "Configura\nmensajes",    "Admin define nombre, concepto y mensaje por color. Kollybry los recuerda para el siguiente mes"),
    ("📲", "Envío\nmasivo",          "Template aprobado por Meta → llega a cada papá en WhatsApp personal con su nombre y monto exacto"),
]

arrow_color = GREEN_MED
step_w = Inches(2.3)
step_h = Inches(3.8)
gap    = Inches(0.28)
start_x = Inches(0.35)
y_top = Inches(1.5)

for i, (icon, title, desc) in enumerate(steps):
    x = start_x + i * (step_w + gap)
    add_rect(s, x, y_top, step_w, step_h, WHITE)
    add_rect(s, x, y_top, step_w, Inches(0.08), GREEN_MED)
    # circle number
    add_rect(s, x + step_w/2 - Inches(0.3), y_top + Inches(0.2), Inches(0.6), Inches(0.6), GREEN_DARK)
    add_text(s, str(i+1), x + step_w/2 - Inches(0.3), y_top+Inches(0.2), Inches(0.6), Inches(0.6),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, icon, x, y_top+Inches(1.0), step_w, Inches(0.6),
             size=26, align=PP_ALIGN.CENTER, color=DARK_TEXT)
    add_text(s, title, x, y_top+Inches(1.6), step_w, Inches(0.7),
             size=13, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_text(s, desc, x+Inches(0.1), y_top+Inches(2.25), step_w-Inches(0.2), Inches(1.45),
             size=10, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
    # arrow (not after last)
    if i < len(steps) - 1:
        ax = x + step_w + Inches(0.02)
        add_rect(s, ax, y_top + step_h/2 - Inches(0.05), gap, Inches(0.1), GREEN_MED)
        add_text(s, "▶", ax - Inches(0.05), y_top + step_h/2 - Inches(0.22), gap + Inches(0.1), Inches(0.44),
                 size=14, color=GREEN_MED, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.35), Inches(5.6), Inches(12.63), Inches(1.45), GREEN_XLIGHT)
add_text(s, "💡  Resultado típico",
         Inches(0.5), Inches(5.65), Inches(3), Inches(0.4),
         size=13, bold=True, color=GREEN_DARK)
add_text(s,
    "300 mensajes personalizados enviados en 2 minutos  •  "
    "Tasa de apertura >95%  •  "
    "Configuración guardada para el próximo mes  •  "
    "Admin no necesita conocimientos técnicos",
    Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.8),
    size=11, color=GRAY_TEXT)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Stack técnico
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Arquitectura y Stack Técnico", "SaaS multi-tenant sobre infraestructura moderna y escalable")

# left — stack table
col_titles = ["Capa", "Tecnología"]
rows_stack = [
    ("Frontend",        "React 18 · Vite · Tailwind CSS · shadcn/ui"),
    ("Backend",         "Node.js 20 · Express · ESM modules"),
    ("Base de datos",   "Supabase PostgreSQL (principal) · Railway PG (CRM)"),
    ("Auth",            "Supabase Auth — JWT · Google OAuth · email/password"),
    ("Mensajería",      "Meta WhatsApp Business API (WABA oficial)"),
    ("Inteligencia artificial", "Anthropic Claude Sonnet 4.6 (mapeo de columnas)"),
    ("Pagos / SPEI",    "Conekta — CLABE única por familia"),
    ("Deploy",          "Railway (backend + frontend)"),
    ("Tiempo real",     "Socket.IO — actualizaciones en vivo"),
]

lx = Inches(0.4)
tw = Inches(6.5)
th = Inches(0.44)
ty = Inches(1.42)
for i, (layer, tech) in enumerate(rows_stack):
    fill = GREEN_XLIGHT if i % 2 == 0 else WHITE
    add_rect(s, lx, ty, Inches(2.0), th, GREEN_DARK if i == -1 else fill)
    add_rect(s, lx + Inches(2.05), ty, Inches(4.4), th, fill)
    c = GREEN_DARK if i % 2 == 0 else DARK_TEXT
    add_text(s, layer, lx+Inches(0.05), ty, Inches(1.95), th, size=10, bold=True, color=GREEN_DARK)
    add_text(s, tech,  lx+Inches(2.1),  ty, Inches(4.3),  th, size=10, color=GRAY_TEXT)
    ty += th

# right — templates Meta
rx = Inches(7.2)
add_rect(s, rx, Inches(1.42), Inches(5.7), Inches(5.5), WHITE)
add_rect(s, rx, Inches(1.42), Inches(5.7), Inches(0.06), GREEN_MED)
add_text(s, "Templates aprobados por Meta",
         rx+Inches(0.12), Inches(1.5), Inches(5.5), Inches(0.38),
         size=13, bold=True, color=GREEN_DARK)

templates = [
    ("kollybry_comunicado_util",      "UTILITY",    "Comunicados — sin límite de frecuencia"),
    ("kollybry_recordatorio_pago",    "UTILITY",    "Recordatorio con monto y link de pago"),
    ("kollybry_aviso_vencido",        "UTILITY",    "Aviso de pago vencido"),
    ("kollybry_confirmacion_pago",    "UTILITY",    "Confirmación de pago recibido"),
    ("kollybry_bienvenida_comunidad", "UTILITY",    "Bienvenida a nuevo contacto"),
    ("kollybry_comunicado_imagen",    "MARKETING",  "Comunicado con imagen adjunta"),
    ("kollybry_recordatorio_spei",    "UTILITY",    "Recordatorio con CLABE SPEI"),
]
ty2 = Inches(2.0)
for tname, cat, desc in templates:
    cat_color = GREEN_DARK if cat == "UTILITY" else ACCENT_TEAL
    add_rect(s, rx+Inches(0.12), ty2, Inches(5.45), Inches(0.48), GREEN_XLIGHT)
    add_text(s, tname, rx+Inches(0.18), ty2+Pt(2), Inches(3.4), Inches(0.44),
             size=9, bold=True, color=GREEN_DARK)
    add_rect(s, rx+Inches(3.65), ty2+Pt(5), Inches(0.85), Inches(0.3), cat_color)
    add_text(s, cat, rx+Inches(3.65), ty2+Pt(5), Inches(0.85), Inches(0.3),
             size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, rx+Inches(0.18), ty2+Inches(0.3), Inches(5.3), Inches(0.22),
             size=9, italic=True, color=GRAY_TEXT)
    ty2 += Inches(0.62)

add_rect(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.7), GREEN_XLIGHT)
add_text(s, "Multi-tenancy: cada colegio/condominio/gym es un tenant con UUID propio. "
            "Todos los datos aislados por tenant_id. Un solo número WhatsApp Business atiende múltiples tenants.",
         Inches(0.55), Inches(6.32), Inches(12.3), Inches(0.64),
         size=10, color=GRAY_TEXT)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Objeciones frecuentes
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Objeciones frecuentes", "Respuestas listas para el equipo de ventas")

qa = [
    ('"Ya tenemos grupo de WhatsApp"',
     "Los grupos mezclan quejas y exponen el celular del maestro/admin. Kollybry envía mensajes institucionales, unidireccionales, con el nombre de la organización como remitente. Sin ruido, sin privacidad expuesta."),
    ('"Usamos correo electrónico"',
     "El correo tiene 20–30% de apertura y tarda 6–12 horas. WhatsApp: 95–98% en < 3 minutos. Para cobros vencidos o avisos urgentes, el correo simplemente no llega a tiempo."),
    ('"¿Los papás/residentes tienen que instalar algo?"',
     "No. Reciben el mensaje en su WhatsApp personal, el que ya usan todos los días. Cero fricción de adopción — no hay que convencer a nadie de descargar una app."),
    ('"¿Es legal mandar mensajes masivos por WhatsApp?"',
     "Sí. Kollybry usa la API oficial de WhatsApp Business de Meta con templates aprobados. El receptor es alumno, residente o socio activo — existe relación previa. Cumple con políticas de Meta y LFPDPPP."),
    ('"¿Qué pasa si alguien responde?"',
     "Las respuestas llegan a la bandeja de WhatsApp Business del administrador, igual que hoy con cualquier WA Business. Kollybry es outbound — los mensajes van del org a su comunidad, no al revés."),
    ('"Ya tenemos una app propia"',
     "La tasa de instalación de apps propias rara vez supera el 40%. El 60% de tus padres, residentes o socios nunca la instaló. WhatsApp ya está instalado en el 100% de sus teléfonos."),
]

cols = 2
card_w = Inches(6.1)
card_h = Inches(1.55)
for i, (q, a) in enumerate(qa):
    col = i % cols
    row = i // cols
    x = Inches(0.4) + col * (card_w + Inches(0.3))
    y = Inches(1.5) + row * (card_h + Inches(0.12))
    add_rect(s, x, y, card_w, card_h, WHITE)
    add_rect(s, x, y, Inches(0.06), card_h, GREEN_MED)
    add_text(s, q, x+Inches(0.12), y+Inches(0.1), card_w-Inches(0.2), Inches(0.38),
             size=11, bold=True, color=GREEN_DARK, italic=True)
    add_text(s, a, x+Inches(0.12), y+Inches(0.5), card_w-Inches(0.2), Inches(1.0),
             size=10, color=GRAY_TEXT)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Cierre
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
add_rect(s, 0, 0, W, H, GREEN_DARK)
add_rect(s, 0, Inches(5.0), W, Inches(2.5), GREEN_MED)

add_text(s, "KOLLYBRY", Inches(1.0), Inches(1.2), Inches(11), Inches(1.4),
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "El mensaje correcto · Al receptor correcto · En el canal correcto",
         Inches(1.0), Inches(2.6), Inches(11), Inches(0.6),
         size=20, color=GREEN_LIGHT, align=PP_ALIGN.CENTER, italic=True)

chips_data = ["95–98% apertura", "< 3 min de entrega", "Cero instalación",
              "Templates Meta aprobados", "IA para cobranza", "Multi-tenant"]
chip_w = Inches(2.0)
chip_h = Inches(0.44)
total_w = len(chips_data) * chip_w + (len(chips_data)-1) * Inches(0.15)
start = (W - total_w) / 2
for i, chip in enumerate(chips_data):
    cx = start + i * (chip_w + Inches(0.15))
    add_rect(s, cx, Inches(3.8), chip_w, chip_h, WHITE)
    add_text(s, chip, cx, Inches(3.8), chip_w, chip_h,
             size=10, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

add_text(s, "nkuvo.com", Inches(1.0), Inches(5.3), Inches(11), Inches(0.5),
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Nkuvo Labs — Querétaro, México · 2026",
         Inches(1.0), Inches(5.85), Inches(11), Inches(0.4),
         size=13, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)
add_text(s, "Construido con WhatsApp Business API oficial de Meta · Stack: React · Node.js · Supabase · Railway · Anthropic AI",
         Inches(1.0), Inches(6.4), Inches(11), Inches(0.4),
         size=10, color=RGBColor(0x77,0xD9,0xB5), align=PP_ALIGN.CENTER, italic=True)

# ── Guardar ────────────────────────────────────────────────────
out = "/Users/alejandrosoria/Colibri-Collections/docs/Kollybry-Pitch-Deck.pptx"
prs.save(out)
print(f"✅ Guardado: {out}")
